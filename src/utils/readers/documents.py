# pyre-ignore-all-errors
"""Readers for document formats: plain text, DOCX, PDF, spreadsheets, presentations.

Each public ``read_X_file`` function accepts either a path (legacy) or an
open binary file-like via the ``fileobj`` keyword. The file-like path is the
SafeDir-friendly entry point: callers open via ``SafeDir.open_for_reader``,
wrap in ``os.fdopen(fd, "rb")``, and hand to the reader. Path-based callers
continue to work unchanged.

The underlying parse logic lives in private ``_parse_X`` helpers so the
fileobj and path branches share a single implementation.
"""

from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import BinaryIO

try:
    import fitz  # PyMuPDF

    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import docx

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text as _rtf_to_text

    STRIPRTF_AVAILABLE = True
except ImportError:
    STRIPRTF_AVAILABLE = False

try:
    import openpyxl

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from loguru import logger

from utils.readers._base import FileReadError, _check_fd_size, _check_file_size


def _parse_text(fileobj: BinaryIO, max_chars: int, label: str) -> str:
    """Decode bytes from *fileobj* as UTF-8 (errors='ignore') and truncate."""
    data = fileobj.read(max_chars * 4)  # *4 to allow worst-case multi-byte chars
    text = data.decode("utf-8", errors="ignore")[:max_chars]
    logger.debug(f"Read {len(text)} characters from {label}")
    return text


def read_text_file(
    file_path: str | Path | None = None,
    max_chars: int = 5000,
    *,
    fileobj: BinaryIO | None = None,
) -> str:
    """Read text content from a plain text file.

    Args:
        file_path: Path to text file (legacy entry point).
        max_chars: Maximum characters to read.
        fileobj: Open binary file-like (SafeDir-friendly entry point). When
            given, ``file_path`` is used only for the log label.

    Returns:
        Text content (UTF-8 decoded, errors ignored, truncated).

    Raises:
        FileReadError: If file cannot be read.
        ValueError: If neither ``file_path`` nor ``fileobj`` is provided.
    """
    if fileobj is not None:
        label = Path(file_path).name if file_path is not None else "<fileobj>"
        try:
            _check_fd_size(fileobj)
            return _parse_text(fileobj, max_chars, label)
        except OSError as e:
            raise FileReadError(f"Failed to read text file {label}: {e}") from e
    if file_path is None:
        raise ValueError("read_text_file requires file_path or fileobj")
    path = Path(file_path)
    _check_file_size(path)
    try:
        with path.open("rb") as f:
            return _parse_text(f, max_chars, path.name)
    except OSError as e:
        raise FileReadError(f"Failed to read text file {path}: {e}") from e


def _parse_docx(fileobj: BinaryIO, label: str) -> str:
    doc = docx.Document(fileobj)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n".join(paragraphs)
    logger.debug(f"Extracted {len(text)} characters from {label}")
    return text


def read_docx_file(
    file_path: str | Path | None = None,
    *,
    fileobj: BinaryIO | None = None,
) -> str:
    """Read text content from a .docx file.

    Args:
        file_path: Path to DOCX file (legacy entry point).
        fileobj: Open binary file-like (SafeDir-friendly entry point).

    Returns:
        Extracted text content.

    Raises:
        FileReadError: If file cannot be read.
        ImportError: If python-docx is not installed.
        ValueError: If neither ``file_path`` nor ``fileobj`` is provided.
    """
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx is not installed. Install with: pip install python-docx")
    if fileobj is not None:
        label = Path(file_path).name if file_path is not None else "<fileobj>"
        try:
            _check_fd_size(fileobj)
            return _parse_docx(fileobj, label)
        except Exception as e:  # Intentional catch-all: python-docx raises library-specific errors
            raise FileReadError(f"Failed to read DOCX file {label}: {e}") from e
    if file_path is None:
        raise ValueError("read_docx_file requires file_path or fileobj")
    path = Path(file_path)
    _check_file_size(path)
    try:
        with path.open("rb") as f:
            return _parse_docx(f, path.name)
    except Exception as e:  # Intentional catch-all: python-docx raises library-specific errors
        raise FileReadError(f"Failed to read DOCX file {path}: {e}") from e


def _extract_pdf_pages(doc: object, max_pages: int, label: str) -> str:
    """Extract text from the first ``max_pages`` of an open PyMuPDF document."""
    num_pages = min(max_pages, len(doc))  # type: ignore[arg-type]
    pages_text = [
        doc.load_page(i).get_text()  # type: ignore[attr-defined]
        for i in range(num_pages)
    ]
    text = "\n".join(pages_text)
    logger.debug(f"Extracted {len(text)} characters from {num_pages} pages of {label}")
    return text


def _parse_pdf_stream(fileobj: BinaryIO, max_pages: int, label: str) -> str:
    """Parse PDF from a file-like by reading bytes and using fitz.open(stream=).

    Used only when the caller already has a fileobj (the SafeDir-friendly
    entry point); the path branch streams from disk via ``fitz.open(path)``
    to avoid loading the whole PDF into memory before ``max_pages`` is
    applied.
    """
    data = fileobj.read()
    with fitz.open(stream=data, filetype="pdf") as doc:
        return _extract_pdf_pages(doc, max_pages, label)


def read_pdf_file(
    file_path: str | Path | None = None,
    max_pages: int = 5,
    *,
    fileobj: BinaryIO | None = None,
) -> str:
    """Read text content from a PDF file.

    Args:
        file_path: Path to PDF file (legacy entry point).
        max_pages: Maximum pages to read.
        fileobj: Open binary file-like (SafeDir-friendly entry point).

    Returns:
        Extracted text content.

    Raises:
        FileReadError: If file cannot be read.
        ImportError: If PyMuPDF is not installed.
        ValueError: If neither ``file_path`` nor ``fileobj`` is provided.
    """
    if not PYMUPDF_AVAILABLE:
        raise ImportError("PyMuPDF is not installed. Install with: pip install PyMuPDF")
    if fileobj is not None:
        label = Path(file_path).name if file_path is not None else "<fileobj>"
        try:
            _check_fd_size(fileobj)
            return _parse_pdf_stream(fileobj, max_pages, label)
        except Exception as e:  # Intentional catch-all: PyMuPDF raises library-specific errors
            raise FileReadError(f"Failed to read PDF file {label}: {e}") from e
    if file_path is None:
        raise ValueError("read_pdf_file requires file_path or fileobj")
    path = Path(file_path)
    _check_file_size(path)
    try:
        with fitz.open(str(path)) as doc:
            return _extract_pdf_pages(doc, max_pages, path.name)
    except Exception as e:  # Intentional catch-all: PyMuPDF raises library-specific errors
        raise FileReadError(f"Failed to read PDF file {path}: {e}") from e


def _parse_rtf(fileobj: BinaryIO, max_chars: int, label: str) -> str:
    raw = fileobj.read().decode("latin-1", errors="ignore")
    text: str = str(_rtf_to_text(raw))
    out = text[:max_chars] if len(text) > max_chars else text
    logger.debug(f"Read {len(out)} characters from RTF {label}")
    return out


def read_rtf_file(
    file_path: str | Path | None = None,
    max_chars: int = 50_000,
    *,
    fileobj: BinaryIO | None = None,
) -> str:
    """Extract plain text from an RTF file using striprtf."""
    if not STRIPRTF_AVAILABLE:
        raise ImportError("striprtf is required for RTF support: pip install striprtf")
    if fileobj is not None:
        label = Path(file_path).name if file_path is not None else "<fileobj>"
        try:
            _check_fd_size(fileobj)
            return _parse_rtf(fileobj, max_chars, label)
        except Exception as exc:
            raise FileReadError(f"Failed to read RTF {label}: {exc}") from exc
    if file_path is None:
        raise ValueError("read_rtf_file requires file_path or fileobj")
    path = Path(file_path)
    _check_file_size(path)
    try:
        with path.open("rb") as f:
            return _parse_rtf(f, max_chars, path.name)
    except Exception as exc:
        raise FileReadError(f"Failed to read RTF {path.name}: {exc}") from exc


def _parse_csv(fileobj: BinaryIO, max_rows: int, label: str) -> str:
    # csv.reader needs text mode; wrap the binary fileobj.
    text_stream = io.TextIOWrapper(fileobj, encoding="utf-8", errors="ignore", newline="")
    rows = []
    reader = csv.reader(text_stream)
    for i, row in enumerate(reader):
        if i >= max_rows:
            break
        rows.append(",".join(row))
    text = "\n".join(rows)
    logger.debug(f"Extracted {len(text)} characters from {len(rows)} rows of {label}")
    return text


def _parse_xlsx(fileobj: BinaryIO, max_rows: int, label: str) -> str:
    wb = openpyxl.load_workbook(fileobj, data_only=True, read_only=True)
    ws = wb.active
    rows = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i >= max_rows:
            break
        row_str = ",".join(str(cell) if cell is not None else "" for cell in row)
        if row_str.strip(","):
            rows.append(row_str)
    text = "\n".join(rows)
    logger.debug(f"Extracted {len(text)} characters from {len(rows)} rows of {label}")
    return text


def _dispatch_spreadsheet(fileobj: BinaryIO, ext: str, max_rows: int, label: str) -> str:
    """Route to ``_parse_csv`` or ``_parse_xlsx`` by extension."""
    if ext == ".csv":
        return _parse_csv(fileobj, max_rows, label)
    if ext in (".xlsx", ".xls"):
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl is not installed. Install with: pip install openpyxl")
        return _parse_xlsx(fileobj, max_rows, label)
    raise FileReadError(f"Unsupported spreadsheet format: {ext}")


def read_spreadsheet_file(
    file_path: str | Path | None = None,
    max_rows: int = 100,
    *,
    fileobj: BinaryIO | None = None,
) -> str:
    """Read content from Excel or CSV file.

    Args:
        file_path: Path to spreadsheet file (legacy entry point, also used to
            detect extension when ``fileobj`` is provided).
        max_rows: Maximum rows to read.
        fileobj: Open binary file-like (SafeDir-friendly entry point). The
            extension is taken from ``file_path`` so the caller must still
            pass the original name; only the I/O is routed through ``fileobj``.

    Returns:
        String representation of data.

    Raises:
        FileReadError: If file cannot be read or extension is unsupported.
        ImportError: If openpyxl is not installed (for Excel files).
        ValueError: If neither ``file_path`` nor ``fileobj`` is provided.
    """
    if file_path is None:
        raise ValueError(
            "read_spreadsheet_file requires file_path (also needed for extension detection "
            "when fileobj is provided)"
        )
    path = Path(file_path)
    ext = path.suffix.lower()
    if fileobj is not None:
        try:
            _check_fd_size(fileobj)
            return _dispatch_spreadsheet(fileobj, ext, max_rows, path.name)
        except (ImportError, FileReadError):
            raise
        except Exception as e:  # Intentional catch-all: openpyxl/csv raise library-specific errors
            raise FileReadError(f"Failed to read spreadsheet file {path.name}: {e}") from e
    _check_file_size(path)
    try:
        with path.open("rb") as f:
            return _dispatch_spreadsheet(f, ext, max_rows, path.name)
    except (ImportError, FileReadError):
        raise
    except Exception as e:  # Intentional catch-all: openpyxl/csv raise library-specific errors
        raise FileReadError(f"Failed to read spreadsheet file {path}: {e}") from e


def _parse_presentation(fileobj: BinaryIO, label: str) -> str:
    prs = Presentation(fileobj)
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"Slide {slide_num}: " + " | ".join(slide_content))
    text = "\n".join(slides_text)
    logger.debug(f"Extracted {len(text)} characters from {len(slides_text)} slides of {label}")
    return text


def read_presentation_file(
    file_path: str | Path | None = None,
    *,
    fileobj: BinaryIO | None = None,
) -> str:
    """Read text content from PowerPoint file.

    Args:
        file_path: Path to PPT/PPTX file (legacy entry point).
        fileobj: Open binary file-like (SafeDir-friendly entry point).

    Returns:
        Extracted text from all slides.

    Raises:
        FileReadError: If file cannot be read.
        ImportError: If python-pptx is not installed.
        ValueError: If neither ``file_path`` nor ``fileobj`` is provided.
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")
    if fileobj is not None:
        label = Path(file_path).name if file_path is not None else "<fileobj>"
        try:
            _check_fd_size(fileobj)
            return _parse_presentation(fileobj, label)
        except Exception as e:  # Intentional catch-all: python-pptx raises library-specific errors
            raise FileReadError(f"Failed to read presentation file {label}: {e}") from e
    if file_path is None:
        raise ValueError("read_presentation_file requires file_path or fileobj")
    path = Path(file_path)
    _check_file_size(path)
    try:
        with path.open("rb") as f:
            return _parse_presentation(f, path.name)
    except Exception as e:  # Intentional catch-all: python-pptx raises library-specific errors
        raise FileReadError(f"Failed to read presentation file {path}: {e}") from e
