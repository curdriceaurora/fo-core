"""Integration tests for utils/readers/documents.py.

Covers:
- read_text_file: basic read, max_chars truncation, missing file raises FileReadError,
  unicode content, empty file
- read_docx_file: ImportError when docx unavailable, real DOCX via python-docx,
  FileReadError on corrupt file
- read_pdf_file: ImportError when PyMuPDF unavailable, real PDF via PyMuPDF,
  multi-page cap, FileReadError on corrupt file
- read_spreadsheet_file: CSV happy path, max_rows truncation, CSV with empty rows,
  XLSX via openpyxl, unsupported extension raises FileReadError,
  ImportError when openpyxl unavailable
- read_presentation_file: ImportError when pptx unavailable, real PPTX via python-pptx,
  empty slide (no text shapes), FileReadError on corrupt file
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

import pytest

from file_organizer.utils.readers._base import FileReadError

pytestmark = pytest.mark.integration


# ---------------------------------------------------------------------------
# read_text_file
# ---------------------------------------------------------------------------


class TestReadTextFile:
    """Integration tests for read_text_file."""

    def test_basic_read(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_text_file

        f = tmp_path / "hello.txt"
        f.write_text("Hello, world!\nSecond line.", encoding="utf-8")

        result = read_text_file(f)

        assert "Hello, world!" in result
        assert "Second line." in result

    def test_max_chars_truncation(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_text_file

        f = tmp_path / "long.txt"
        f.write_text("A" * 1000, encoding="utf-8")

        result = read_text_file(f, max_chars=100)

        assert len(result) == 100
        assert result == "A" * 100

    def test_missing_file_raises_file_read_error(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_text_file

        missing = tmp_path / "nonexistent.txt"

        with pytest.raises(FileReadError, match="Failed to read text file"):
            read_text_file(missing)

    def test_unicode_content(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_text_file

        f = tmp_path / "unicode.txt"
        f.write_text("こんにちは — Héllo Wörld", encoding="utf-8")

        result = read_text_file(f)

        assert "Héllo" in result

    def test_empty_file(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_text_file

        f = tmp_path / "empty.txt"
        f.write_text("", encoding="utf-8")

        result = read_text_file(f)

        assert result == ""

    def test_returns_string(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_text_file

        f = tmp_path / "data.txt"
        f.write_text("content", encoding="utf-8")

        result = read_text_file(f)

        assert isinstance(result, str)
        assert result == "content"


# ---------------------------------------------------------------------------
# read_docx_file
# ---------------------------------------------------------------------------


class TestReadDocxFile:
    """Integration tests for read_docx_file using python-docx."""

    @pytest.fixture(autouse=True)
    def _require_docx(self) -> None:
        pytest.importorskip("docx")

    def test_real_docx_extraction(self, tmp_path: Path) -> None:
        import docx as python_docx

        from file_organizer.utils.readers.documents import read_docx_file

        doc_path = tmp_path / "test.docx"
        doc = python_docx.Document()
        doc.add_paragraph("First paragraph content.")
        doc.add_paragraph("Second paragraph content.")
        doc.save(str(doc_path))

        result = read_docx_file(doc_path)

        assert "First paragraph content." in result
        assert "Second paragraph content." in result

    def test_empty_paragraphs_skipped(self, tmp_path: Path) -> None:
        import docx as python_docx

        from file_organizer.utils.readers.documents import read_docx_file

        doc_path = tmp_path / "sparse.docx"
        doc = python_docx.Document()
        doc.add_paragraph("")
        doc.add_paragraph("   ")
        doc.add_paragraph("Real content here.")
        doc.save(str(doc_path))

        result = read_docx_file(doc_path)

        assert "Real content here." in result
        assert result.strip() == "Real content here."

    def test_import_error_when_docx_unavailable(self) -> None:
        from file_organizer.utils.readers.documents import read_docx_file

        with patch("file_organizer.utils.readers.documents.DOCX_AVAILABLE", False):
            with pytest.raises(ImportError, match="python-docx is not installed"):
                read_docx_file("test.docx")

    def test_corrupt_docx_raises_file_read_error(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_docx_file

        bad_file = tmp_path / "corrupt.docx"
        bad_file.write_bytes(b"not a docx file at all")

        with pytest.raises(FileReadError, match="Failed to read DOCX file"):
            read_docx_file(bad_file)

    def test_returns_joined_paragraphs(self, tmp_path: Path) -> None:
        import docx as python_docx

        from file_organizer.utils.readers.documents import read_docx_file

        doc_path = tmp_path / "multi.docx"
        doc = python_docx.Document()
        doc.add_paragraph("Alpha")
        doc.add_paragraph("Beta")
        doc.add_paragraph("Gamma")
        doc.save(str(doc_path))

        result = read_docx_file(doc_path)

        assert result == "Alpha\nBeta\nGamma"


# ---------------------------------------------------------------------------
# read_pdf_file
# ---------------------------------------------------------------------------


class TestReadPdfFile:
    """Integration tests for read_pdf_file using PyMuPDF (fitz)."""

    @pytest.fixture(autouse=True)
    def _require_fitz(self) -> None:
        pytest.importorskip("fitz")

    def test_import_error_when_pymupdf_unavailable(self) -> None:
        from file_organizer.utils.readers.documents import read_pdf_file

        with patch("file_organizer.utils.readers.documents.PYMUPDF_AVAILABLE", False):
            with pytest.raises(ImportError, match="PyMuPDF is not installed"):
                read_pdf_file("test.pdf")

    def test_corrupt_pdf_raises_file_read_error(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_pdf_file

        bad_pdf = tmp_path / "corrupt.pdf"
        bad_pdf.write_bytes(b"not a pdf at all")

        with pytest.raises(FileReadError, match="Failed to read PDF file"):
            read_pdf_file(bad_pdf)

    def test_max_pages_respected(self, tmp_path: Path) -> None:
        import fitz

        from file_organizer.utils.readers.documents import read_pdf_file

        doc_path = tmp_path / "multi.pdf"
        doc = fitz.open()
        for i in range(4):
            page = doc.new_page()
            page.insert_text((50, 100), f"Page {i + 1} text content")
        doc.save(str(doc_path))
        doc.close()

        result = read_pdf_file(doc_path, max_pages=2)

        assert "Page 1" in result
        assert "Page 2" in result
        assert "Page 4" not in result

    def test_single_page_pdf(self, tmp_path: Path) -> None:
        import fitz

        from file_organizer.utils.readers.documents import read_pdf_file

        doc_path = tmp_path / "single.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 100), "Single page content here.")
        doc.save(str(doc_path))
        doc.close()

        result = read_pdf_file(doc_path)

        assert "Single page content here." in result

    def test_returns_string(self, tmp_path: Path) -> None:
        import fitz

        from file_organizer.utils.readers.documents import read_pdf_file

        doc_path = tmp_path / "check.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 100), "Check content.")
        doc.save(str(doc_path))
        doc.close()

        result = read_pdf_file(doc_path)

        assert isinstance(result, str)
        assert len(result) > 0


# ---------------------------------------------------------------------------
# read_spreadsheet_file
# ---------------------------------------------------------------------------


class TestReadSpreadsheetFile:
    """Integration tests for read_spreadsheet_file."""

    def test_csv_basic_read(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        csv_file = tmp_path / "data.csv"
        csv_file.write_text("Name,Age\nAlice,30\nBob,25", encoding="utf-8")

        result = read_spreadsheet_file(csv_file)

        assert "Name,Age" in result
        assert "Alice,30" in result
        assert "Bob,25" in result

    def test_csv_max_rows_truncation(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        csv_file = tmp_path / "big.csv"
        lines = ["col1,col2"] + [f"row{i},{i}" for i in range(200)]
        csv_file.write_text("\n".join(lines), encoding="utf-8")

        result = read_spreadsheet_file(csv_file, max_rows=5)

        rows = result.strip().split("\n")
        assert len(rows) == 5
        assert "col1,col2" in rows[0]

    def test_csv_returns_string(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        csv_file = tmp_path / "simple.csv"
        csv_file.write_text("a,b,c\n1,2,3", encoding="utf-8")

        result = read_spreadsheet_file(csv_file)

        assert isinstance(result, str)
        assert "a,b,c" in result

    def test_csv_single_column(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        csv_file = tmp_path / "single_col.csv"
        csv_file.write_text("value\nfoo\nbar\nbaz", encoding="utf-8")

        result = read_spreadsheet_file(csv_file)

        assert "value" in result
        assert "foo" in result

    def test_unsupported_format_raises_file_read_error(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        bad_file = tmp_path / "data.ods"
        bad_file.touch()

        with pytest.raises(FileReadError, match="Unsupported spreadsheet format"):
            read_spreadsheet_file(bad_file)

    def test_xlsx_basic_read(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        openpyxl = pytest.importorskip("openpyxl")

        xlsx_path = tmp_path / "sheet.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Product"
        ws["B1"] = "Price"
        ws["A2"] = "Widget"
        ws["B2"] = 9.99
        wb.save(str(xlsx_path))

        result = read_spreadsheet_file(xlsx_path)

        assert "Product" in result
        assert "Widget" in result

    def test_xlsx_max_rows(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        openpyxl = pytest.importorskip("openpyxl")

        xlsx_path = tmp_path / "large.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        for i in range(1, 201):
            ws[f"A{i}"] = f"row{i}"
        wb.save(str(xlsx_path))

        result = read_spreadsheet_file(xlsx_path, max_rows=10)

        rows = result.strip().split("\n")
        assert len(rows) == 10
        assert "row1" in result
        assert "row11" not in result

    def test_xlsx_empty_cells_handled(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        openpyxl = pytest.importorskip("openpyxl")

        xlsx_path = tmp_path / "sparse.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "filled"
        ws["B1"] = None
        ws["C1"] = "also filled"
        wb.save(str(xlsx_path))

        result = read_spreadsheet_file(xlsx_path)

        assert "filled" in result
        assert "also filled" in result

    def test_openpyxl_unavailable_raises_import_error(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_spreadsheet_file

        xlsx_path = tmp_path / "test.xlsx"
        xlsx_path.touch()

        with patch("file_organizer.utils.readers.documents.OPENPYXL_AVAILABLE", False):
            with pytest.raises(ImportError, match="openpyxl is not installed"):
                read_spreadsheet_file(xlsx_path)


# ---------------------------------------------------------------------------
# read_presentation_file
# ---------------------------------------------------------------------------


class TestReadPresentationFile:
    """Integration tests for read_presentation_file using python-pptx."""

    @pytest.fixture(autouse=True)
    def _require_pptx(self) -> None:
        pytest.importorskip("pptx")

    def test_import_error_when_pptx_unavailable(self) -> None:
        from file_organizer.utils.readers.documents import read_presentation_file

        with patch("file_organizer.utils.readers.documents.PPTX_AVAILABLE", False):
            with pytest.raises(ImportError, match="python-pptx is not installed"):
                read_presentation_file("test.pptx")

    def test_real_pptx_extraction(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Pt

        from file_organizer.utils.readers.documents import read_presentation_file

        pptx_path = tmp_path / "slides.pptx"
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        slide1 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide1.shapes.add_textbox(0, 0, Pt(200), Pt(50))
        txBox.text_frame.text = "Slide one title"

        slide2 = prs.slides.add_slide(blank_slide_layout)
        txBox2 = slide2.shapes.add_textbox(0, 0, Pt(200), Pt(50))
        txBox2.text_frame.text = "Slide two content"

        prs.save(str(pptx_path))

        result = read_presentation_file(pptx_path)

        assert "Slide 1" in result
        assert "Slide one title" in result
        assert "Slide 2" in result
        assert "Slide two content" in result

    def test_empty_slides_not_included(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Pt

        from file_organizer.utils.readers.documents import read_presentation_file

        pptx_path = tmp_path / "mixed.pptx"
        prs = Presentation()

        blank_layout = prs.slide_layouts[6]

        slide1 = prs.slides.add_slide(blank_layout)
        txBox = slide1.shapes.add_textbox(0, 0, Pt(200), Pt(50))
        txBox.text_frame.text = "Has content"

        prs.slides.add_slide(blank_layout)

        prs.save(str(pptx_path))

        result = read_presentation_file(pptx_path)

        assert "Has content" in result
        assert "Slide 2" not in result

    def test_corrupt_pptx_raises_file_read_error(self, tmp_path: Path) -> None:
        from file_organizer.utils.readers.documents import read_presentation_file

        bad_file = tmp_path / "corrupt.pptx"
        bad_file.write_bytes(b"not a pptx file")

        with pytest.raises(FileReadError, match="Failed to read presentation file"):
            read_presentation_file(bad_file)

    def test_returns_newline_joined_slides(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Pt

        from file_organizer.utils.readers.documents import read_presentation_file

        pptx_path = tmp_path / "multi_slide.pptx"
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]

        for i in range(1, 4):
            slide = prs.slides.add_slide(blank_layout)
            tb = slide.shapes.add_textbox(0, 0, Pt(200), Pt(50))
            tb.text_frame.text = f"Content for slide {i}"

        prs.save(str(pptx_path))

        result = read_presentation_file(pptx_path)

        lines = result.strip().split("\n")
        assert len(lines) == 3
        assert all(line.startswith("Slide ") for line in lines)
