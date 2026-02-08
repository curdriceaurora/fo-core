"""File reading utilities for various file types."""

from pathlib import Path

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    EBOOKLIB_AVAILABLE = True
except ImportError:
    EBOOKLIB_AVAILABLE = False

try:
    import ezdxf
    EZDXF_AVAILABLE = True
except ImportError:
    EZDXF_AVAILABLE = False

# Archive format support (built-in)
import zipfile
import tarfile

try:
    import py7zr
    PY7ZR_AVAILABLE = True
except ImportError:
    PY7ZR_AVAILABLE = False

try:
    import rarfile
    RARFILE_AVAILABLE = True
except ImportError:
    RARFILE_AVAILABLE = False

# Scientific format support
try:
    import h5py
    H5PY_AVAILABLE = True
except ImportError:
    H5PY_AVAILABLE = False

try:
    import netCDF4
    NETCDF4_AVAILABLE = True
except ImportError:
    NETCDF4_AVAILABLE = False

try:
    from scipy.io import loadmat
    SCIPY_AVAILABLE = True
except ImportError:
    SCIPY_AVAILABLE = False

from loguru import logger


class FileReadError(Exception):
    """Exception raised when file reading fails."""
    pass


def read_text_file(file_path: str | Path, max_chars: int = 5000) -> str:
    """Read text content from a plain text file.

    Args:
        file_path: Path to text file
        max_chars: Maximum characters to read

    Returns:
        Text content

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    try:
        with open(file_path, encoding='utf-8', errors='ignore') as f:
            text = f.read(max_chars)
        logger.debug(f"Read {len(text)} characters from {file_path.name}")
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read text file {file_path}: {e}") from e


def read_docx_file(file_path: str | Path) -> str:
    """Read text content from a .docx file.

    Args:
        file_path: Path to DOCX file

    Returns:
        Extracted text content

    Raises:
        FileReadError: If file cannot be read
        ImportError: If python-docx is not installed
    """
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx is not installed. Install with: pip install python-docx")

    file_path = Path(file_path)
    try:
        doc = docx.Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        text = '\n'.join(paragraphs)
        logger.debug(f"Extracted {len(text)} characters from {file_path.name}")
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read DOCX file {file_path}: {e}") from e


def read_pdf_file(file_path: str | Path, max_pages: int = 5) -> str:
    """Read text content from a PDF file.

    Args:
        file_path: Path to PDF file
        max_pages: Maximum pages to read

    Returns:
        Extracted text content

    Raises:
        FileReadError: If file cannot be read
        ImportError: If PyMuPDF is not installed
    """
    if not PYMUPDF_AVAILABLE:
        raise ImportError("PyMuPDF is not installed. Install with: pip install PyMuPDF")

    file_path = Path(file_path)
    try:
        doc = fitz.open(file_path)
        num_pages = min(max_pages, len(doc))

        pages_text = []
        for page_num in range(num_pages):
            page = doc.load_page(page_num)
            pages_text.append(page.get_text())

        text = '\n'.join(pages_text)
        doc.close()

        logger.debug(
            f"Extracted {len(text)} characters from {num_pages} pages of {file_path.name}"
        )
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read PDF file {file_path}: {e}") from e


def read_spreadsheet_file(file_path: str | Path, max_rows: int = 100) -> str:
    """Read content from Excel or CSV file.

    Args:
        file_path: Path to spreadsheet file
        max_rows: Maximum rows to read

    Returns:
        String representation of data

    Raises:
        FileReadError: If file cannot be read
        ImportError: If pandas is not installed
    """
    if not PANDAS_AVAILABLE:
        raise ImportError("pandas is not installed. Install with: pip install pandas openpyxl")

    file_path = Path(file_path)
    try:
        # Determine file type and read
        if file_path.suffix.lower() == '.csv':
            df = pd.read_csv(file_path, nrows=max_rows)
        elif file_path.suffix.lower() in ('.xlsx', '.xls'):
            df = pd.read_excel(file_path, nrows=max_rows)
        else:
            raise ValueError(f"Unsupported spreadsheet format: {file_path.suffix}")

        # Convert to string, limiting size
        text = df.to_string(max_rows=max_rows)

        logger.debug(
            f"Extracted {len(text)} characters from {len(df)} rows of {file_path.name}"
        )
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read spreadsheet file {file_path}: {e}") from e


def read_presentation_file(file_path: str | Path) -> str:
    """Read text content from PowerPoint file.

    Args:
        file_path: Path to PPT/PPTX file

    Returns:
        Extracted text from all slides

    Raises:
        FileReadError: If file cannot be read
        ImportError: If python-pptx is not installed
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")

    file_path = Path(file_path)
    try:
        prs = Presentation(file_path)

        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

            if slide_content:
                slides_text.append(f"Slide {slide_num}: " + " | ".join(slide_content))

        text = '\n'.join(slides_text)
        logger.debug(
            f"Extracted {len(text)} characters from {len(slides_text)} slides of {file_path.name}"
        )
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read presentation file {file_path}: {e}") from e


def read_ebook_file(file_path: str | Path, max_chars: int = 10000) -> str:
    """Read text content from ebook file (EPUB only for now).

    Args:
        file_path: Path to ebook file
        max_chars: Maximum characters to extract

    Returns:
        Extracted text content

    Raises:
        FileReadError: If file cannot be read
        ImportError: If ebooklib is not installed
    """
    if not EBOOKLIB_AVAILABLE:
        raise ImportError("ebooklib is not installed. Install with: pip install ebooklib")

    file_path = Path(file_path)

    # Only support EPUB for now
    if file_path.suffix.lower() != '.epub':
        raise ValueError(f"Unsupported ebook format: {file_path.suffix}. Only .epub supported.")

    try:
        book = epub.read_epub(file_path)

        text_parts = []
        total_chars = 0

        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                content = item.get_content().decode('utf-8', errors='ignore')
                # Basic HTML stripping (simple approach)
                import re
                content = re.sub(r'<[^>]+>', ' ', content)
                content = re.sub(r'\s+', ' ', content).strip()

                if content:
                    text_parts.append(content)
                    total_chars += len(content)

                    if total_chars >= max_chars:
                        break

        text = ' '.join(text_parts)[:max_chars]

        logger.debug(f"Extracted {len(text)} characters from ebook {file_path.name}")
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read ebook file {file_path}: {e}") from e


def read_zip_file(file_path: str | Path, max_files: int = 50) -> str:
    """Read contents and metadata from a ZIP archive.

    Args:
        file_path: Path to ZIP file
        max_files: Maximum number of files to list

    Returns:
        String with archive metadata and file listing

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            info_list = zf.infolist()[:max_files]

            # Calculate statistics
            total_files = len(zf.infolist())
            total_compressed = sum(info.compress_size for info in zf.infolist())
            total_uncompressed = sum(info.file_size for info in zf.infolist())
            compression_ratio = (1 - total_compressed / total_uncompressed) * 100 if total_uncompressed > 0 else 0

            # Check for encryption
            encrypted = any(info.flag_bits & 0x1 for info in zf.infolist())

            # Build metadata string
            lines = [
                f"ZIP Archive: {file_path.name}",
                f"Total files: {total_files}",
                f"Compressed size: {total_compressed / 1024:.2f} KB",
                f"Uncompressed size: {total_uncompressed / 1024:.2f} KB",
                f"Compression ratio: {compression_ratio:.1f}%",
                f"Encrypted: {'Yes' if encrypted else 'No'}",
                "\nFiles (first {}):" .format(min(max_files, total_files)),
            ]

            # List files
            for info in info_list:
                size_kb = info.file_size / 1024
                compressed_kb = info.compress_size / 1024
                lines.append(
                    f"  - {info.filename} ({size_kb:.2f} KB → {compressed_kb:.2f} KB)"
                )

            if total_files > max_files:
                lines.append(f"  ... and {total_files - max_files} more files")

            text = '\n'.join(lines)
            logger.debug(f"Extracted metadata from ZIP archive {file_path.name} ({total_files} files)")
            return text

    except Exception as e:
        raise FileReadError(f"Failed to read ZIP file {file_path}: {e}") from e


def read_7z_file(file_path: str | Path, max_files: int = 50) -> str:
    """Read contents and metadata from a 7Z archive.

    Args:
        file_path: Path to 7Z file
        max_files: Maximum number of files to list

    Returns:
        String with archive metadata and file listing

    Raises:
        FileReadError: If file cannot be read
        ImportError: If py7zr is not installed
    """
    if not PY7ZR_AVAILABLE:
        raise ImportError("py7zr is not installed. Install with: pip install py7zr")

    file_path = Path(file_path)
    try:
        with py7zr.SevenZipFile(file_path, 'r') as archive:
            all_files = archive.list()

            # Calculate statistics
            total_files = len(all_files)
            total_compressed = sum(f.compressed for f in all_files)
            total_uncompressed = sum(f.uncompressed for f in all_files)
            compression_ratio = (1 - total_compressed / total_uncompressed) * 100 if total_uncompressed > 0 else 0

            # Check for encryption
            encrypted = archive.password_protected if hasattr(archive, 'password_protected') else False

            # Build metadata string
            lines = [
                f"7Z Archive: {file_path.name}",
                f"Total files: {total_files}",
                f"Compressed size: {total_compressed / 1024:.2f} KB",
                f"Uncompressed size: {total_uncompressed / 1024:.2f} KB",
                f"Compression ratio: {compression_ratio:.1f}%",
                f"Encrypted: {'Yes' if encrypted else 'No'}",
                "\nFiles (first {}):" .format(min(max_files, total_files)),
            ]

            # List files
            for idx, file_info in enumerate(all_files[:max_files]):
                size_kb = file_info.uncompressed / 1024
                compressed_kb = file_info.compressed / 1024
                lines.append(
                    f"  - {file_info.filename} ({size_kb:.2f} KB → {compressed_kb:.2f} KB)"
                )

            if total_files > max_files:
                lines.append(f"  ... and {total_files - max_files} more files")

            text = '\n'.join(lines)
            logger.debug(f"Extracted metadata from 7Z archive {file_path.name} ({total_files} files)")
            return text

    except Exception as e:
        raise FileReadError(f"Failed to read 7Z file {file_path}: {e}") from e


def read_tar_file(file_path: str | Path, max_files: int = 50) -> str:
    """Read contents and metadata from a TAR/GZ/BZ2 archive.

    Args:
        file_path: Path to TAR file (.tar, .tar.gz, .tgz, .tar.bz2)
        max_files: Maximum number of files to list

    Returns:
        String with archive metadata and file listing

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    try:
        with tarfile.open(file_path, 'r:*') as tf:
            members = tf.getmembers()

            # Calculate statistics
            total_files = len([m for m in members if m.isfile()])
            total_dirs = len([m for m in members if m.isdir()])
            total_size = sum(m.size for m in members if m.isfile())

            # Determine compression type
            compression_type = "None"
            if file_path.suffix in ('.gz', '.tgz'):
                compression_type = "GZ"
            elif file_path.suffix in ('.bz2', '.tbz2'):
                compression_type = "BZ2"
            elif file_path.suffix == '.xz':
                compression_type = "XZ"

            # Build metadata string
            lines = [
                f"TAR Archive: {file_path.name}",
                f"Compression: {compression_type}",
                f"Total files: {total_files}",
                f"Total directories: {total_dirs}",
                f"Total size: {total_size / 1024:.2f} KB",
                f"\nFiles (first {min(max_files, total_files)}):",
            ]

            # List files (skip directories)
            file_members = [m for m in members if m.isfile()][:max_files]
            for member in file_members:
                size_kb = member.size / 1024
                lines.append(f"  - {member.name} ({size_kb:.2f} KB)")

            if total_files > max_files:
                lines.append(f"  ... and {total_files - max_files} more files")

            text = '\n'.join(lines)
            logger.debug(f"Extracted metadata from TAR archive {file_path.name} ({total_files} files)")
            return text

    except Exception as e:
        raise FileReadError(f"Failed to read TAR file {file_path}: {e}") from e


def read_rar_file(file_path: str | Path, max_files: int = 50) -> str:
    """Read contents and metadata from a RAR archive.

    Args:
        file_path: Path to RAR file
        max_files: Maximum number of files to list

    Returns:
        String with archive metadata and file listing

    Raises:
        FileReadError: If file cannot be read
        ImportError: If rarfile is not installed
    """
    if not RARFILE_AVAILABLE:
        raise ImportError(
            "rarfile is not installed. Install with: pip install rarfile\n"
            "Note: RAR support also requires unrar command-line tool to be installed."
        )

    file_path = Path(file_path)
    try:
        with rarfile.RarFile(file_path, 'r') as rf:
            info_list = rf.infolist()

            # Calculate statistics
            total_files = len(info_list)
            total_compressed = sum(info.compress_size for info in info_list)
            total_uncompressed = sum(info.file_size for info in info_list)
            compression_ratio = (1 - total_compressed / total_uncompressed) * 100 if total_uncompressed > 0 else 0

            # Check for encryption
            encrypted = rf.needs_password()

            # Build metadata string
            lines = [
                f"RAR Archive: {file_path.name}",
                f"Total files: {total_files}",
                f"Compressed size: {total_compressed / 1024:.2f} KB",
                f"Uncompressed size: {total_uncompressed / 1024:.2f} KB",
                f"Compression ratio: {compression_ratio:.1f}%",
                f"Encrypted: {'Yes' if encrypted else 'No'}",
                "\nFiles (first {}):" .format(min(max_files, total_files)),
            ]

            # List files
            for info in info_list[:max_files]:
                size_kb = info.file_size / 1024
                compressed_kb = info.compress_size / 1024
                lines.append(
                    f"  - {info.filename} ({size_kb:.2f} KB → {compressed_kb:.2f} KB)"
                )

            if total_files > max_files:
                lines.append(f"  ... and {total_files - max_files} more files")

            text = '\n'.join(lines)
            logger.debug(f"Extracted metadata from RAR archive {file_path.name} ({total_files} files)")
            return text

    except Exception as e:
        raise FileReadError(f"Failed to read RAR file {file_path}: {e}") from e


def read_hdf5_file(file_path: str | Path, max_datasets: int = 20) -> str:
    """Read metadata and structure from an HDF5 file.

    Args:
        file_path: Path to HDF5 file
        max_datasets: Maximum number of datasets to list

    Returns:
        String with HDF5 structure and metadata

    Raises:
        FileReadError: If file cannot be read
        ImportError: If h5py is not installed
    """
    if not H5PY_AVAILABLE:
        raise ImportError("h5py is not installed. Install with: pip install h5py")

    file_path = Path(file_path)
    try:
        with h5py.File(file_path, 'r') as hf:
            lines = [
                f"HDF5 File: {file_path.name}",
                f"Total groups: {len(list(hf.keys()))}",
                "\nStructure:",
            ]

            dataset_count = 0

            def visit_item(name: str, obj: h5py.Dataset | h5py.Group) -> None:
                nonlocal dataset_count
                if dataset_count >= max_datasets:
                    return

                if isinstance(obj, h5py.Dataset):
                    shape_str = 'x'.join(map(str, obj.shape))
                    size_kb = obj.nbytes / 1024
                    lines.append(
                        f"  Dataset: {name} [{obj.dtype}] {shape_str} ({size_kb:.2f} KB)"
                    )

                    # List attributes
                    if obj.attrs:
                        for attr_name, attr_value in list(obj.attrs.items())[:3]:
                            lines.append(f"    - {attr_name}: {attr_value}")

                    dataset_count += 1
                elif isinstance(obj, h5py.Group):
                    lines.append(f"  Group: {name}/")

            hf.visititems(visit_item)

            if dataset_count >= max_datasets:
                lines.append(f"  ... (showing first {max_datasets} datasets)")

            text = '\n'.join(lines)
            logger.debug(f"Extracted metadata from HDF5 file {file_path.name}")
            return text

    except Exception as e:
        raise FileReadError(f"Failed to read HDF5 file {file_path}: {e}") from e


def read_netcdf_file(file_path: str | Path) -> str:
    """Read metadata and structure from a NetCDF file.

    Args:
        file_path: Path to NetCDF file

    Returns:
        String with NetCDF structure and metadata

    Raises:
        FileReadError: If file cannot be read
        ImportError: If netCDF4 is not installed
    """
    if not NETCDF4_AVAILABLE:
        raise ImportError("netCDF4 is not installed. Install with: pip install netCDF4")

    file_path = Path(file_path)
    try:
        with netCDF4.Dataset(file_path, 'r') as nc:
            lines = [
                f"NetCDF File: {file_path.name}",
                f"Format: {nc.data_model}",
                "\nDimensions:",
            ]

            # List dimensions
            for dim_name, dim in nc.dimensions.items():
                size = len(dim) if not dim.isunlimited() else "unlimited"
                lines.append(f"  - {dim_name}: {size}")

            lines.append("\nVariables:")

            # List variables (first 20)
            for idx, (var_name, var) in enumerate(list(nc.variables.items())[:20]):
                shape_str = 'x'.join(str(var.shape[i]) for i in range(len(var.shape)))
                lines.append(f"  - {var_name} ({var.dtype}): {shape_str}")

                # Show some attributes
                if hasattr(var, 'units'):
                    lines.append(f"      units: {var.units}")
                if hasattr(var, 'long_name'):
                    lines.append(f"      long_name: {var.long_name}")

            if len(nc.variables) > 20:
                lines.append(f"  ... and {len(nc.variables) - 20} more variables")

            # Global attributes
            if nc.ncattrs():
                lines.append("\nGlobal Attributes:")
                for attr_name in list(nc.ncattrs())[:10]:
                    attr_value = nc.getncattr(attr_name)
                    lines.append(f"  - {attr_name}: {attr_value}")

            text = '\n'.join(lines)
            logger.debug(f"Extracted metadata from NetCDF file {file_path.name}")
            return text

    except Exception as e:
        raise FileReadError(f"Failed to read NetCDF file {file_path}: {e}") from e


def read_mat_file(file_path: str | Path) -> str:
    """Read metadata and structure from a MATLAB .mat file.

    Args:
        file_path: Path to MAT file

    Returns:
        String with MAT file structure and metadata

    Raises:
        FileReadError: If file cannot be read
        ImportError: If scipy is not installed
    """
    if not SCIPY_AVAILABLE:
        raise ImportError("scipy is not installed. Install with: pip install scipy")

    file_path = Path(file_path)
    try:
        # Load mat file
        mat_contents = loadmat(file_path, struct_as_record=False, squeeze_me=True)

        lines = [
            f"MATLAB File: {file_path.name}",
            "\nVariables:",
        ]

        # Filter out metadata variables
        var_names = [k for k in mat_contents.keys() if not k.startswith('__')]

        for var_name in var_names[:30]:  # Limit to first 30 variables
            var = mat_contents[var_name]

            # Get type and shape info
            var_type = type(var).__name__
            if hasattr(var, 'shape'):
                shape_str = 'x'.join(map(str, var.shape))
                lines.append(f"  - {var_name} ({var_type}): {shape_str}")
            else:
                lines.append(f"  - {var_name} ({var_type})")

        if len(var_names) > 30:
            lines.append(f"  ... and {len(var_names) - 30} more variables")

        text = '\n'.join(lines)
        logger.debug(f"Extracted metadata from MAT file {file_path.name}")
        return text

    except Exception as e:
        raise FileReadError(f"Failed to read MAT file {file_path}: {e}") from e


def read_file(file_path: str | Path, **kwargs) -> str | None:
    """Read content from any supported file type.

    Auto-detects file type and uses appropriate reader.

    Args:
        file_path: Path to file
        **kwargs: Additional arguments passed to specific readers

    Returns:
        Extracted text content, or None if unsupported

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)

    # Check for compound extensions (e.g., .tar.gz)
    name_lower = file_path.name.lower()
    ext = file_path.suffix.lower()

    # Handle compound extensions for archives
    if name_lower.endswith('.tar.gz') or name_lower.endswith('.tar.bz2') or name_lower.endswith('.tar.xz'):
        compound_ext = '.' + '.'.join(file_path.name.split('.')[-2:]).lower()
    else:
        compound_ext = ext

    readers = {
        # Document formats
        ('.txt', '.md'): read_text_file,
        ('.docx',): read_docx_file,  # Note: .doc (old binary format) is NOT supported
        ('.pdf',): read_pdf_file,
        ('.csv', '.xlsx', '.xls'): read_spreadsheet_file,
        ('.ppt', '.pptx'): read_presentation_file,
        ('.epub',): read_ebook_file,
        # Archive formats
        ('.zip',): read_zip_file,
        ('.7z',): read_7z_file,
        ('.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2', '.tar.xz'): read_tar_file,
        ('.rar',): read_rar_file,
        # Scientific formats
        ('.hdf5', '.h5', '.hdf'): read_hdf5_file,
        ('.nc', '.nc4', '.netcdf'): read_netcdf_file,
        ('.mat',): read_mat_file,
        # CAD formats
        ('.dxf', '.dwg', '.step', '.stp', '.iges', '.igs'): read_cad_file,
    }

    # Try compound extension first, then fall back to simple extension
    for check_ext in [compound_ext, ext]:
        for extensions, reader in readers.items():
            if check_ext in extensions:
                try:
                    return reader(file_path, **kwargs)
                except Exception as e:
                    logger.error(f"Error reading {file_path.name}: {e}")
                    raise

    logger.warning(f"Unsupported file type: {ext}")
    return None


def read_dxf_file(file_path: str | Path, max_layers: int = 20) -> str:
    """Read metadata and content from a DXF CAD file.

    Args:
        file_path: Path to DXF file
        max_layers: Maximum number of layers to list

    Returns:
        Extracted metadata and layer information

    Raises:
        FileReadError: If file cannot be read
        ImportError: If ezdxf is not installed
    """
    if not EZDXF_AVAILABLE:
        raise ImportError("ezdxf is not installed. Install with: pip install ezdxf")

    file_path = Path(file_path)
    try:
        doc = ezdxf.readfile(file_path)

        # Extract document information
        metadata_parts = []

        # Header variables
        if hasattr(doc, 'header'):
            metadata_parts.append("=== DXF Document Metadata ===")

            # Try to get common header variables safely
            try:
                title = doc.header.get('$TITLE', 'Untitled')
                if title:
                    metadata_parts.append(f"Title: {title}")
            except:
                pass

            try:
                author = doc.header.get('$AUTHOR', 'Unknown')
                if author:
                    metadata_parts.append(f"Author: {author}")
            except:
                pass

        # DXF version
        metadata_parts.append(f"DXF Version: {doc.dxfversion}")

        # Layer information
        if doc.layers:
            layer_count = len(doc.layers)
            metadata_parts.append(f"\n=== Layers ({layer_count} total) ===")

            for idx, layer in enumerate(doc.layers):
                if idx >= max_layers:
                    metadata_parts.append(f"... and {layer_count - max_layers} more layers")
                    break

                layer_info = f"Layer: {layer.dxf.name}"
                if hasattr(layer.dxf, 'color'):
                    layer_info += f" (Color: {layer.dxf.color})"
                metadata_parts.append(layer_info)

        # Entity statistics
        modelspace = doc.modelspace()
        entity_types: dict[str, int] = {}

        for entity in modelspace:
            entity_type = entity.dxftype()
            entity_types[entity_type] = entity_types.get(entity_type, 0) + 1

        if entity_types:
            metadata_parts.append(f"\n=== Entities ===")
            metadata_parts.append(f"Total entities: {sum(entity_types.values())}")

            # List entity types
            for entity_type, count in sorted(entity_types.items()):
                metadata_parts.append(f"  {entity_type}: {count}")

        # Blocks
        if doc.blocks:
            block_count = len([b for b in doc.blocks if not b.name.startswith('*')])
            if block_count > 0:
                metadata_parts.append(f"\n=== Blocks ===")
                metadata_parts.append(f"Block definitions: {block_count}")

        text = '\n'.join(metadata_parts)
        logger.debug(f"Extracted {len(text)} characters from DXF file {file_path.name}")
        return text

    except Exception as e:
        raise FileReadError(f"Failed to read DXF file {file_path}: {e}") from e


def read_dwg_file(file_path: str | Path) -> str:
    """Read metadata from a DWG CAD file.

    Note: DWG is a proprietary format. This function attempts to read it using ezdxf,
    which has limited DWG support. For full DWG support, consider using ODA File Converter
    to convert DWG to DXF first.

    Args:
        file_path: Path to DWG file

    Returns:
        Extracted metadata or basic file information

    Raises:
        FileReadError: If file cannot be read
        ImportError: If ezdxf is not installed
    """
    if not EZDXF_AVAILABLE:
        raise ImportError("ezdxf is not installed. Install with: pip install ezdxf")

    file_path = Path(file_path)
    try:
        # Try to read with ezdxf (limited support)
        doc = ezdxf.readfile(file_path)
        return read_dxf_file(file_path)  # Process as DXF

    except Exception as e:
        # If ezdxf can't read it, provide basic file info
        logger.warning(f"Could not parse DWG file with ezdxf: {e}")

        metadata_parts = [
            "=== DWG File Information ===",
            f"File: {file_path.name}",
            f"Size: {file_path.stat().st_size / 1024:.2f} KB",
            "",
            "Note: Full DWG parsing requires additional tools.",
            "Consider using ODA File Converter to convert DWG to DXF for better support."
        ]

        return '\n'.join(metadata_parts)


def read_step_file(file_path: str | Path, max_lines: int = 100) -> str:
    """Read metadata from a STEP (.step, .stp) CAD file.

    STEP files are ISO 10303 standard format for 3D CAD data exchange.
    This function extracts basic header information.

    Args:
        file_path: Path to STEP file
        max_lines: Maximum lines to read from header

    Returns:
        Extracted header information

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read(10000)  # Read first 10KB

        metadata_parts = ["=== STEP File Information ==="]
        metadata_parts.append(f"File: {file_path.name}")
        metadata_parts.append(f"Size: {file_path.stat().st_size / 1024:.2f} KB")

        # Extract header section (between HEADER; and ENDSEC;)
        if 'HEADER;' in content and 'ENDSEC;' in content:
            header_start = content.find('HEADER;')
            header_end = content.find('ENDSEC;', header_start)
            header = content[header_start:header_end]

            metadata_parts.append("\n=== Header Information ===")

            # Extract file description
            if 'FILE_DESCRIPTION' in header:
                desc_start = header.find('FILE_DESCRIPTION')
                desc_end = header.find(');', desc_start)
                if desc_end > desc_start:
                    desc = header[desc_start:desc_end + 2]
                    metadata_parts.append(desc.strip())

            # Extract file name
            if 'FILE_NAME' in header:
                name_start = header.find('FILE_NAME')
                name_end = header.find(');', name_start)
                if name_end > name_start:
                    name_info = header[name_start:name_end + 2]
                    metadata_parts.append(name_info.strip())

            # Extract schema
            if 'FILE_SCHEMA' in header:
                schema_start = header.find('FILE_SCHEMA')
                schema_end = header.find(');', schema_start)
                if schema_end > schema_start:
                    schema = header[schema_start:schema_end + 2]
                    metadata_parts.append(schema.strip())

        # Count entities in DATA section
        if 'DATA;' in content:
            # Count lines starting with # (entities)
            entity_count = content.count('\n#')
            metadata_parts.append(f"\nApproximate entity count: {entity_count}")

        text = '\n'.join(metadata_parts)
        logger.debug(f"Extracted {len(text)} characters from STEP file {file_path.name}")
        return text

    except Exception as e:
        raise FileReadError(f"Failed to read STEP file {file_path}: {e}") from e


def read_iges_file(file_path: str | Path, max_lines: int = 50) -> str:
    """Read metadata from an IGES (.iges, .igs) CAD file.

    IGES (Initial Graphics Exchange Specification) is a vendor-neutral file format
    for 3D CAD data exchange.

    Args:
        file_path: Path to IGES file
        max_lines: Maximum lines to read from header

    Returns:
        Extracted header information

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            lines = [f.readline() for _ in range(max_lines)]

        metadata_parts = ["=== IGES File Information ==="]
        metadata_parts.append(f"File: {file_path.name}")
        metadata_parts.append(f"Size: {file_path.stat().st_size / 1024:.2f} KB")

        # IGES files have structured sections marked in column 73
        # S = Start, G = Global, D = Directory Entry, P = Parameter Data, T = Terminate

        start_section = []
        global_section = []

        for line in lines:
            if len(line) >= 73:
                section_type = line[72]
                content = line[:72].strip()

                if section_type == 'S' and content:
                    start_section.append(content)
                elif section_type == 'G' and content:
                    global_section.append(content)

        # Display start section (usually contains file description)
        if start_section:
            metadata_parts.append("\n=== Start Section ===")
            metadata_parts.extend(start_section[:10])  # First 10 lines

        # Display global section (contains parameters)
        if global_section:
            metadata_parts.append("\n=== Global Parameters ===")
            # Join and split by commas to get parameters
            global_params = ''.join(global_section)
            params = global_params.split(',')[:5]  # First 5 parameters

            # Parameter meanings (typical order)
            param_names = [
                "Parameter delimiter",
                "Record delimiter",
                "Product ID from sender",
                "File name",
                "Native system ID"
            ]

            for i, (name, value) in enumerate(zip(param_names, params)):
                if value.strip():
                    metadata_parts.append(f"{name}: {value.strip()}")

        # Count entities
        entity_count = sum(1 for line in lines if len(line) >= 73 and line[72] == 'D')
        if entity_count > 0:
            metadata_parts.append(f"\nDirectory entries found: {entity_count}")

        text = '\n'.join(metadata_parts)
        logger.debug(f"Extracted {len(text)} characters from IGES file {file_path.name}")
        return text

    except Exception as e:
        raise FileReadError(f"Failed to read IGES file {file_path}: {e}") from e


def read_cad_file(file_path: str | Path, **kwargs) -> str:
    """Read content from any CAD file format.

    Auto-detects CAD file type and uses appropriate reader.

    Args:
        file_path: Path to CAD file
        **kwargs: Additional arguments passed to specific readers

    Returns:
        Extracted metadata and content

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    cad_readers = {
        '.dxf': read_dxf_file,
        '.dwg': read_dwg_file,
        '.step': read_step_file,
        '.stp': read_step_file,
        '.iges': read_iges_file,
        '.igs': read_iges_file,
    }

    reader = cad_readers.get(ext)
    if reader:
        return reader(file_path, **kwargs)
    else:
        raise ValueError(f"Unsupported CAD file format: {ext}")
